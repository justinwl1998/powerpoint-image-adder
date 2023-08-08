from pptx import Presentation
from PIL import Image


# Adds images and scales them based on the size of the
# placeholder in the presentation.
def add_image(slide, placeholder_id, img_path):
    placeholder = slide.placeholders[placeholder_id]
    im = Image.open(img_path)
    width, height = im.size

    placeholder.height = height
    placeholder.width = width

    placeholder = placeholder.insert_picture(img_path)

    image_ratio = width / height
    placeholder_ratio = placeholder.width / placeholder.height
    ratio_difference = placeholder_ratio - image_ratio

    if ratio_difference > 0:
        difference_on_each_side = ratio_difference / 2
        placeholder.crop_left = -difference_on_each_side
        placeholder.crop_right = -difference_on_each_side
    else:
        difference_on_each_side = -ratio_difference / 2
        placeholder.crop_left = -difference_on_each_side
        placeholder.crop_right = -difference_on_each_side

# Loops through every slide and inserts an image in each placeholder
# that takes an image.
def populateSlides(images, presentation, progress, counter_label):    
    prs = Presentation(presentation)
    
    resultCode = 0
    index = 0
    imgIndex = 0
    hits = 0
    breakOut = False

    while not breakOut:
        try:
            # Case where there are not enough slides with placeholders to hold the images
            if index >= len(prs.slides):
                resultCode = 2
                break
            slide = prs.slides[index]
            for shape in slide.placeholders:
                # check if placeholder is already occupied
                if type(shape).__name__ == "PlaceholderPicture":
                    continue
                
                if type(shape).__name__ == "PicturePlaceholder":
                    # Case where there are not enough images for placeholders
                    if imgIndex >= len(images):
                        breakOut = True
                        resultCode = 1
                        break
                    hits += 1
                    add_image(slide, shape.placeholder_format.idx, images[imgIndex])
                    imgIndex += 1
                    counter_label.set(str(hits) + " image(s) added")
            index += 1
        except:
            # error code case
            return 4

    progress['value'] = 100
    prs.save(presentation)
    return 0


    
