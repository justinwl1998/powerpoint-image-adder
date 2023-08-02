from pptx import Presentation

prs = Presentation()


title_slide_layout = prs.slide_masters[0].slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
a, b, c, d = slide.placeholders

title.text = "Hello, world!"
a.text = 'placeholder text 1'
b.text = 'placeholder text 2'
c.text = 'placeholder text 3'
d.text = 'placeholder text 4'
#subtitle.text = "Working with a new library is hard."


prs.save('test.pptx')

#Todo:

# Audomate adding images to an existing powerpoint presentation

# ask for directory where images to use in presentation
# ask for the powerpoint presentation to add to
# go through each slide after the title and add four images maximum
#  to each slide

# learn how add_picture() can be used

# make a powerpoint slide style where four images can be added to it

