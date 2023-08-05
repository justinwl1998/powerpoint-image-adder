import os
import copy
from threading import Thread
from tkinter import filedialog as fd
from tkinter import messagebox
from tkinter import ttk
from tkinter import *
from pptx import Presentation
from pptx.util import Inches
from PIL import Image
from helpers import add_image, populateSlides

def checkInputs():
    print(imageList)
    if len(root.splitlist(entry_text.get())) > 0 and template_text.get() != "":
        print("Both conditions are met!")
        B3['state'] = NORMAL
    else:
        print("Both conditions are not met.")
        B3['state'] = DISABLED

def imageSelectCallBack():
    selected = fd.askopenfilename(
        parent=root,
        multiple=True,
        filetypes=[('Image files', ['.jpg', '.png'])])

    imageList = []
    for img in selected:
        print(img)
        imageList.append(img)

    print(imageList)
        
    #print(list(selected))
    entry_text.set(selected)
    #print(entry_text.get())
    num_label.set(str(len(root.splitlist(selected))) + " images")
    
    checkInputs()
    return

def templateSelectCallBack():
    selected = fd.askopenfilename(title="Select template presentation file",
                              filetypes=[('Presentation files', ['.pptx'])])
    template_text.set(selected)
    checkInputs()
    return


def populateCallback():
    B1['state'] = "disabled"
    B2['state'] = "disabled"
    B3['state'] = "disabled"


    for img in entry_text.get():
        print(img)
    print(5/0)
    
    populateSlides(list(entry_text.get()), template_text.get())
    msg_box = messagebox.showinfo(message="Wrote to presentation")
    B1['state'] = "normal"
    B2['state'] = "normal"
    B3['state'] = "normal"
    return

root = Tk()
imageList = []
root.title('Powerpoint Mass Image Adder')
root.resizable(False, False)
root.geometry('500x300')

L1 = Label(root, text="Images directory", justify="left", anchor="w")
L1.grid(column=0, row=3, sticky=W, pady=(15,0), padx=(30,0))

entry_text = StringVar()
num_label = StringVar()
num_label.set("0 images")
imgCountLabel = Label(root, textvariable=num_label, justify="left", anchor="w")
imgCountLabel.grid(column=0, row=4, sticky=W, padx=(30,0))
B1 = Button(root, text="...", command=imageSelectCallBack, width=8)
B1.grid(column=1, row=4, sticky=W)

#Entry for empty presentation file to populate with images
L2 = Label(root, text="Template presentation", justify="left", anchor="w")
L2.grid(column=0, row=6, sticky=W, pady=(15,0), padx=(30,0))

template_text = StringVar()
template_text.set('D:/coding/powerpoint-image-adder/template.pptx')
E2 = Entry(root, width=55, bd=5, state=DISABLED, textvariable=template_text)
E2.grid(column=0, row=7, sticky=W, padx=(30,0))
B2 = Button(root, text="...", command=templateSelectCallBack, width=8)
B2.grid(column=1, row=7, sticky=W)

#Progress bar eventually

B3 = Button(root, text="Read", command=populateCallback, state=DISABLED)
B3.place(relx=0.5, rely=0.8, anchor=CENTER, height=32, width=150)
root.mainloop()

##if os.path.exists('./test.pptx'):
##    os.remove('test.pptx')
##
##path = "template.pptx"
##
##templatePres = Presentation(path)
##prs = Presentation()
##
##template = templatePres.slides[0]
##
##for shape in template.placeholders:
##    print('%d %s' % (shape.placeholder_format.idx, shape.name))
##
##    if 'Picture Placeholder' in shape.name:
##        add_image(template, shape.placeholder_format.idx)
##
##
##    
##templatePres.save("test.pptx")
#os.startfile("test.pptx")
#prs.save('test.pptx')

#Todo:

# Audomate adding images to an existing powerpoint presentation

# ask for directory where images to use in presentation
# ask for the powerpoint presentation to add to
# go through each slide after the title and add four images maximum
#  to each slide


